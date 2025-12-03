"""Utilities for parsing slide decks (PPT/PPTX) into text with slide numbers."""
from typing import List, Tuple

from pptx import Presentation

from .pdf_utils import chunk_text_with_pages


def extract_text_with_slides(path: str) -> List[Tuple[str, int]]:
    """
    Extract text content from each slide in a PowerPoint file.

    Returns a list of (text, slide_number) tuples.
    """
    prs = Presentation(path)
    slides: List[Tuple[str, int]] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        texts: List[str] = []
        for shape in slide.shapes:
            if not hasattr(shape, "text"):
                continue
            text = (shape.text or "").strip()
            if text:
                texts.append(text)

        if texts:
            slides.append(("\n".join(texts), slide_idx))

    return slides


def chunk_text_with_slides(
    slide_texts: List[Tuple[str, int]],
    max_chars: int = 1500,
    overlap: int = 200,
) -> List[Tuple[str, int]]:
    """
    Chunk slide text while preserving slide numbers.

    Uses the same chunking strategy as PDFs, treating slide numbers
    analogously to page numbers.
    """
    return chunk_text_with_pages(slide_texts, max_chars=max_chars, overlap=overlap)


